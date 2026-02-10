"""
Slide Code Examples - Professional design reference patterns for python-pptx slide creation.

These are reference examples to adapt for your content needs.
Also useful for debugging code errors.

Design Philosophy:
- Bold dominant colors (60-70% coverage), not white backgrounds
- Named color palettes, not generic blue
- Font pairings (serif headers + sans-serif body)
- Visual elements in every slide (shapes, icon circles, accent bars)
- Proper size contrast (titles 36pt+, body 14-16pt)
- 0.5"+ margins, 0.3-0.5" spacing between elements
- No accent lines under titles

Categories:
- text_layout: Icon rows, accent bars, structured text with visual hierarchy
- number_highlight: Hero stats, metric cards with bold typography
- grid_layout: Comparison cards, feature columns with color coding
- image_text: Half-bleed layouts, full background overlays
- visual_emphasis: Quote callouts, process steps with visual motifs
- design_reference: Color palettes and font pairing reference (text, not code)
"""

SLIDE_EXAMPLES = {
    "text_layout": {
        "description": "Icon rows, accent bars, structured text with visual hierarchy",
        "when_to_use": "Multiple text items, lists, structured content with professional styling",
        "examples": [
            {
                "name": "icon_text_rows",
                "code": '''
# Icon circles with bold headers and descriptions - Midnight Executive palette
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Dark background
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = RGBColor(0x1E, 0x27, 0x61)

# Title - Georgia serif header
title = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11), Inches(0.9))
tf = title.text_frame
p = tf.paragraphs[0]
p.text = "Key Findings"
p.font.size = Pt(40)
p.font.bold = True
p.font.name = "Georgia"
p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

# Icon + text rows
items = [
    ("01", "Revenue Growth", "Year-over-year revenue increased by 25%, driven by expansion into new markets and improved retention rates."),
    ("02", "Customer Satisfaction", "Net Promoter Score reached 72, placing us in the top quartile of our industry benchmark group."),
    ("03", "Operational Efficiency", "Process automation reduced turnaround time by 40% while maintaining quality standards across all departments."),
]

for i, (icon_text, heading, description) in enumerate(items):
    y = 1.8 + i * 1.7

    # Icon circle - accent color
    circle = slide.shapes.add_shape(9, Inches(0.8), Inches(y), Inches(0.7), Inches(0.7))
    circle.fill.solid()
    circle.fill.fore_color.rgb = RGBColor(0x40, 0x8E, 0xC6)
    circle.line.fill.background()

    # Icon number inside circle
    icon = slide.shapes.add_textbox(Inches(0.8), Inches(y + 0.1), Inches(0.7), Inches(0.5))
    p = icon.text_frame.paragraphs[0]
    p.text = icon_text
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.name = "Calibri"
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.CENTER

    # Bold header - Calibri
    header = slide.shapes.add_textbox(Inches(1.8), Inches(y), Inches(10), Inches(0.5))
    p = header.text_frame.paragraphs[0]
    p.text = heading
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.name = "Calibri"
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # Description text
    desc = slide.shapes.add_textbox(Inches(1.8), Inches(y + 0.5), Inches(10), Inches(0.9))
    tf = desc.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = description
    p.font.size = Pt(14)
    p.font.name = "Calibri"
    p.font.color.rgb = RGBColor(0xB0, 0xC4, 0xDE)
    p.line_spacing = 1.4
'''
            },
            {
                "name": "left_accent_bar",
                "code": '''
# Thick left accent bar with text hierarchy - Teal Trust palette
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Dark background
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = RGBColor(0x0A, 0x1A, 0x2A)

# Thick left accent bar
accent_bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.35), Inches(7.5))
accent_bar.fill.solid()
accent_bar.fill.fore_color.rgb = RGBColor(0x02, 0x80, 0x90)
accent_bar.line.fill.background()

# Title - large with breathing room
title = slide.shapes.add_textbox(Inches(1.0), Inches(0.8), Inches(11), Inches(1.0))
tf = title.text_frame
p = tf.paragraphs[0]
p.text = "Strategic Direction"
p.font.size = Pt(44)
p.font.bold = True
p.font.name = "Georgia"
p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

# Subtitle with breathing room
subtitle = slide.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11), Inches(0.6))
p = subtitle.text_frame.paragraphs[0]
p.text = "Modernizing our approach for sustainable growth"
p.font.size = Pt(20)
p.font.name = "Calibri"
p.font.color.rgb = RGBColor(0x02, 0x80, 0x90)

# Thin separator line
sep = slide.shapes.add_shape(1, Inches(1.0), Inches(2.9), Inches(3), Inches(0.03))
sep.fill.solid()
sep.fill.fore_color.rgb = RGBColor(0x02, 0x80, 0x90)
sep.line.fill.background()

# Body text with generous line spacing
body = slide.shapes.add_textbox(Inches(1.0), Inches(3.3), Inches(10.5), Inches(3.5))
tf = body.text_frame
tf.word_wrap = True

paragraphs = [
    "Our three-year transformation focuses on cloud-native architecture, AI-driven automation, and customer experience redesign.",
    "Phase 1 targets infrastructure modernization with a projected 30% cost reduction. Phase 2 introduces intelligent workflows that will streamline operations across all business units.",
    "Early pilots have shown promising results with 2.3x improvement in processing speed and 45% reduction in manual intervention."
]

for i, text in enumerate(paragraphs):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(15)
    p.font.name = "Calibri"
    p.font.color.rgb = RGBColor(0xCC, 0xDD, 0xEE)
    p.line_spacing = 1.6
    p.space_after = Pt(14)
'''
            }
        ]
    },

    "number_highlight": {
        "description": "Hero stats and metric cards with bold typography",
        "when_to_use": "Key metrics, statistics, performance data with visual impact",
        "examples": [
            {
                "name": "hero_stat_dark",
                "code": '''
# Massive hero stat on dark background - Forest & Moss palette
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Full dark background
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = RGBColor(0x2C, 0x5F, 0x2D)

# Subtle top accent bar
accent = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(0.06))
accent.fill.solid()
accent.fill.fore_color.rgb = RGBColor(0x97, 0xBC, 0x62)
accent.line.fill.background()

# Context label above number
label_above = slide.shapes.add_textbox(Inches(0), Inches(1.5), prs.slide_width, Inches(0.6))
p = label_above.text_frame.paragraphs[0]
p.text = "YEAR OVER YEAR"
p.font.size = Pt(16)
p.font.bold = True
p.font.name = "Calibri"
p.font.color.rgb = RGBColor(0x97, 0xBC, 0x62)
p.alignment = PP_ALIGN.CENTER
p.font.letter_spacing = Pt(3)

# Massive stat number - 120pt
number = slide.shapes.add_textbox(Inches(0), Inches(2.0), prs.slide_width, Inches(2.5))
tf = number.text_frame
p = tf.paragraphs[0]
p.text = "147%"
p.font.size = Pt(120)
p.font.bold = True
p.font.name = "Arial Black"
p.font.color.rgb = RGBColor(0x97, 0xBC, 0x62)
p.alignment = PP_ALIGN.CENTER

# Label below number
label = slide.shapes.add_textbox(Inches(0), Inches(4.6), prs.slide_width, Inches(0.7))
p = label.text_frame.paragraphs[0]
p.text = "Revenue Growth"
p.font.size = Pt(28)
p.font.name = "Georgia"
p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
p.alignment = PP_ALIGN.CENTER

# Supporting context
context = slide.shapes.add_textbox(Inches(2), Inches(5.5), Inches(9.333), Inches(0.8))
tf = context.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Exceeding our target of 100% growth, driven by enterprise expansion and improved retention"
p.font.size = Pt(14)
p.font.name = "Calibri"
p.font.color.rgb = RGBColor(0xCC, 0xDD, 0xCC)
p.alignment = PP_ALIGN.CENTER
'''
            },
            {
                "name": "metric_cards",
                "code": '''
# Metric cards with colored top accent bars - Berry & Cream palette
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Light warm background
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = RGBColor(0xEC, 0xE2, 0xD0)

# Title
title = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11), Inches(0.9))
p = title.text_frame.paragraphs[0]
p.text = "Q4 Performance Metrics"
p.font.size = Pt(36)
p.font.bold = True
p.font.name = "Georgia"
p.font.color.rgb = RGBColor(0x6D, 0x2E, 0x46)

# KPI data
metrics = [
    ("$2.4M", "Revenue", "+18%"),
    ("1,250", "New Customers", "+32%"),
    ("98.5%", "Uptime SLA", "+0.3%"),
]

card_width = 3.5
start_x = 0.8
spacing = 0.5

for i, (value, label, change) in enumerate(metrics):
    x = start_x + i * (card_width + spacing)

    # Card background (white)
    card_bg = slide.shapes.add_shape(1, Inches(x), Inches(2.0), Inches(card_width), Inches(4.0))
    card_bg.fill.solid()
    card_bg.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    card_bg.line.fill.background()

    # Colored top accent bar on card
    top_bar = slide.shapes.add_shape(1, Inches(x), Inches(2.0), Inches(card_width), Inches(0.08))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = RGBColor(0x6D, 0x2E, 0x46)
    top_bar.line.fill.background()

    # Big number
    val_box = slide.shapes.add_textbox(Inches(x + 0.3), Inches(2.6), Inches(card_width - 0.6), Inches(1.2))
    p = val_box.text_frame.paragraphs[0]
    p.text = value
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.name = "Arial Black"
    p.font.color.rgb = RGBColor(0x6D, 0x2E, 0x46)
    p.alignment = PP_ALIGN.CENTER

    # Label
    lbl_box = slide.shapes.add_textbox(Inches(x + 0.3), Inches(3.8), Inches(card_width - 0.6), Inches(0.5))
    p = lbl_box.text_frame.paragraphs[0]
    p.text = label
    p.font.size = Pt(16)
    p.font.name = "Calibri"
    p.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
    p.alignment = PP_ALIGN.CENTER

    # Change indicator
    chg_box = slide.shapes.add_textbox(Inches(x + 0.3), Inches(4.6), Inches(card_width - 0.6), Inches(0.5))
    p = chg_box.text_frame.paragraphs[0]
    p.text = change
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.name = "Calibri"
    p.font.color.rgb = RGBColor(0x2E, 0x8B, 0x57)
    p.alignment = PP_ALIGN.CENTER
'''
            }
        ]
    },

    "grid_layout": {
        "description": "Comparison cards and feature columns with color coding",
        "when_to_use": "Before/after comparisons, feature lists, multi-column content",
        "examples": [
            {
                "name": "comparison_cards",
                "code": '''
# Before/After comparison cards - Coral Energy palette
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Dark background
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x2E)

# Title
title = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11), Inches(0.9))
p = title.text_frame.paragraphs[0]
p.text = "Transformation Results"
p.font.size = Pt(38)
p.font.bold = True
p.font.name = "Georgia"
p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

columns = [
    {
        "title": "Before",
        "color": RGBColor(0xFF, 0x6F, 0x61),
        "card_bg": RGBColor(0x2A, 0x1A, 0x1A),
        "items": ["Manual data entry processes", "3-day average turnaround", "12% error rate in reporting", "Siloed team communication"]
    },
    {
        "title": "After",
        "color": RGBColor(0x00, 0xB4, 0xAA),
        "card_bg": RGBColor(0x1A, 0x2A, 0x2A),
        "items": ["Fully automated workflows", "Same-day delivery", "99.7% accuracy achieved", "Unified collaboration platform"]
    }
]

card_width = 5.5
start_x = 0.8
spacing = 0.7

for i, col in enumerate(columns):
    x = start_x + i * (card_width + spacing)

    # Card background
    card = slide.shapes.add_shape(1, Inches(x), Inches(1.8), Inches(card_width), Inches(5.0))
    card.fill.solid()
    card.fill.fore_color.rgb = col["card_bg"]
    card.line.fill.background()

    # Column header with color accent
    header_bar = slide.shapes.add_shape(1, Inches(x), Inches(1.8), Inches(card_width), Inches(0.08))
    header_bar.fill.solid()
    header_bar.fill.fore_color.rgb = col["color"]
    header_bar.line.fill.background()

    header = slide.shapes.add_textbox(Inches(x + 0.4), Inches(2.1), Inches(card_width - 0.8), Inches(0.6))
    p = header.text_frame.paragraphs[0]
    p.text = col["title"]
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.name = "Calibri"
    p.font.color.rgb = col["color"]

    # Items with dot indicators
    for j, item in enumerate(col["items"]):
        item_y = 3.0 + j * 0.9

        # Dot indicator
        dot = slide.shapes.add_shape(9, Inches(x + 0.4), Inches(item_y + 0.08), Inches(0.15), Inches(0.15))
        dot.fill.solid()
        dot.fill.fore_color.rgb = col["color"]
        dot.line.fill.background()

        # Item text
        txt = slide.shapes.add_textbox(Inches(x + 0.8), Inches(item_y), Inches(card_width - 1.2), Inches(0.5))
        tf = txt.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = item
        p.font.size = Pt(15)
        p.font.name = "Calibri"
        p.font.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
'''
            },
            {
                "name": "three_column_feature",
                "code": '''
# 3-column features with icon circles - Ocean Gradient palette
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Deep dark background
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = RGBColor(0x06, 0x5A, 0x82)

# Title centered
title = slide.shapes.add_textbox(Inches(0), Inches(0.5), prs.slide_width, Inches(0.9))
p = title.text_frame.paragraphs[0]
p.text = "Our Core Services"
p.font.size = Pt(40)
p.font.bold = True
p.font.name = "Georgia"
p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
p.alignment = PP_ALIGN.CENTER

features = [
    {"icon": "01", "title": "Consulting", "desc": "Expert guidance for digital transformation strategy and execution. We help organizations navigate complex technology decisions."},
    {"icon": "02", "title": "Development", "desc": "Custom software solutions built with modern architectures. From MVPs to enterprise-scale platforms tailored to your needs."},
    {"icon": "03", "title": "Support", "desc": "24/7 technical assistance with dedicated account management. Proactive monitoring ensures maximum uptime and performance."}
]

col_width = 3.5
start_x = 0.9
spacing = 0.5

for i, feat in enumerate(features):
    x = start_x + i * (col_width + spacing)

    # Icon circle
    circle = slide.shapes.add_shape(9, Inches(x + col_width/2 - 0.4), Inches(2.0), Inches(0.8), Inches(0.8))
    circle.fill.solid()
    circle.fill.fore_color.rgb = RGBColor(0x1B, 0x9A, 0xAA)
    circle.line.fill.background()

    # Icon text
    icon_txt = slide.shapes.add_textbox(Inches(x + col_width/2 - 0.4), Inches(2.1), Inches(0.8), Inches(0.7))
    p = icon_txt.text_frame.paragraphs[0]
    p.text = feat["icon"]
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.name = "Calibri"
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.CENTER

    # Feature title
    ft = slide.shapes.add_textbox(Inches(x), Inches(3.1), Inches(col_width), Inches(0.6))
    p = ft.text_frame.paragraphs[0]
    p.text = feat["title"]
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.name = "Calibri"
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.CENTER

    # Description - left aligned
    desc = slide.shapes.add_textbox(Inches(x + 0.2), Inches(3.8), Inches(col_width - 0.4), Inches(2.5))
    tf = desc.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = feat["desc"]
    p.font.size = Pt(14)
    p.font.name = "Calibri"
    p.font.color.rgb = RGBColor(0xBB, 0xDD, 0xEE)
    p.line_spacing = 1.5
'''
            }
        ]
    },

    "image_text": {
        "description": "Half-bleed layouts and full background overlays",
        "when_to_use": "Product showcase, visual storytelling, impact slides",
        "examples": [
            {
                "name": "half_bleed_layout",
                "code": '''
# Colored left half + text right half - Warm Terracotta palette
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os

# Light background for right side
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = RGBColor(0xFA, 0xF5, 0xF0)

# Left half - colored placeholder or image
image_files = [f for f in os.listdir('.') if f.endswith(('.png', '.jpg', '.jpeg'))]
if image_files:
    slide.shapes.add_picture(image_files[0], Inches(0), Inches(0), width=Inches(6.5), height=Inches(7.5))
else:
    # Warm terracotta colored placeholder
    placeholder = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(6.5), Inches(7.5))
    placeholder.fill.solid()
    placeholder.fill.fore_color.rgb = RGBColor(0xC4, 0x6A, 0x4E)
    placeholder.line.fill.background()

    # Placeholder label
    ph_label = slide.shapes.add_textbox(Inches(1.5), Inches(3.2), Inches(3.5), Inches(1))
    p = ph_label.text_frame.paragraphs[0]
    p.text = "IMAGE"
    p.font.size = Pt(32)
    p.font.name = "Calibri"
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    from pptx.enum.text import PP_ALIGN
    p.alignment = PP_ALIGN.CENTER

# Right side text content with generous margins
title = slide.shapes.add_textbox(Inches(7.2), Inches(1.5), Inches(5.3), Inches(1.0))
tf = title.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Crafted with Purpose"
p.font.size = Pt(34)
p.font.bold = True
p.font.name = "Georgia"
p.font.color.rgb = RGBColor(0xC4, 0x6A, 0x4E)

# Thin accent line
accent = slide.shapes.add_shape(1, Inches(7.2), Inches(2.8), Inches(2), Inches(0.04))
accent.fill.solid()
accent.fill.fore_color.rgb = RGBColor(0xC4, 0x6A, 0x4E)
accent.line.fill.background()

# Description
desc = slide.shapes.add_textbox(Inches(7.2), Inches(3.2), Inches(5.3), Inches(3.5))
tf = desc.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Every detail has been thoughtfully designed to deliver an exceptional experience. Our approach combines timeless craftsmanship with modern innovation."
p.font.size = Pt(15)
p.font.name = "Calibri"
p.font.color.rgb = RGBColor(0x4A, 0x3A, 0x30)
p.line_spacing = 1.6
p.space_after = Pt(12)

p2 = tf.add_paragraph()
p2.text = "From concept to delivery, we ensure quality at every step of the journey."
p2.font.size = Pt(15)
p2.font.name = "Calibri"
p2.font.color.rgb = RGBColor(0x4A, 0x3A, 0x30)
p2.line_spacing = 1.6
'''
            },
            {
                "name": "full_bg_overlay",
                "code": '''
# Dark overlay on full background, centered large text - Charcoal Minimal palette
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# Dark charcoal base background
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = RGBColor(0x1C, 0x1C, 0x1E)

# If image available, add it as background
image_files = [f for f in os.listdir('.') if f.endswith(('.png', '.jpg', '.jpeg'))]
if image_files:
    slide.shapes.add_picture(image_files[0], Inches(0), Inches(0), width=Inches(13.333), height=Inches(7.5))
    # Dark overlay rectangle
    overlay = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = RGBColor(0x1C, 0x1C, 0x1E)
    overlay.line.fill.background()

# Top accent line
accent = slide.shapes.add_shape(1, Inches(5.5), Inches(1.8), Inches(2.333), Inches(0.04))
accent.fill.solid()
accent.fill.fore_color.rgb = RGBColor(0xE8, 0xE8, 0xE8)
accent.line.fill.background()

# Large centered title
title = slide.shapes.add_textbox(Inches(1), Inches(2.3), Inches(11.333), Inches(1.5))
p = title.text_frame.paragraphs[0]
p.text = "Innovation Starts Here"
p.font.size = Pt(52)
p.font.bold = True
p.font.name = "Georgia"
p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
p.alignment = PP_ALIGN.CENTER

# Subtitle
subtitle = slide.shapes.add_textbox(Inches(2), Inches(4.0), Inches(9.333), Inches(1))
tf = subtitle.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Building the future through bold ideas, rigorous execution, and unwavering commitment to excellence"
p.font.size = Pt(18)
p.font.name = "Calibri"
p.font.color.rgb = RGBColor(0xAA, 0xAA, 0xAA)
p.alignment = PP_ALIGN.CENTER
p.line_spacing = 1.5

# Bottom accent line
accent2 = slide.shapes.add_shape(1, Inches(5.5), Inches(5.5), Inches(2.333), Inches(0.04))
accent2.fill.solid()
accent2.fill.fore_color.rgb = RGBColor(0xE8, 0xE8, 0xE8)
accent2.line.fill.background()
'''
            }
        ]
    },

    "visual_emphasis": {
        "description": "Quote callouts and process steps with visual motifs",
        "when_to_use": "Key takeaways, important quotes, step-by-step processes",
        "examples": [
            {
                "name": "quote_callout",
                "code": '''
# Large quote with thick vertical accent bar - Cherry Bold palette
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Dark background
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = RGBColor(0x15, 0x0E, 0x11)

# Thick vertical accent bar (left side)
accent = slide.shapes.add_shape(1, Inches(1.0), Inches(1.5), Inches(0.12), Inches(4.5))
accent.fill.solid()
accent.fill.fore_color.rgb = RGBColor(0x99, 0x00, 0x11)
accent.line.fill.background()

# Large quotation mark
quote_mark = slide.shapes.add_textbox(Inches(1.5), Inches(1.2), Inches(2), Inches(1.5))
p = quote_mark.text_frame.paragraphs[0]
p.text = "\u201C"
p.font.size = Pt(80)
p.font.name = "Georgia"
p.font.color.rgb = RGBColor(0x99, 0x00, 0x11)

# Quote text
quote = slide.shapes.add_textbox(Inches(1.5), Inches(2.2), Inches(10), Inches(2.5))
tf = quote.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "The companies that will thrive are those willing to reimagine everything about how they create and deliver value."
p.font.size = Pt(28)
p.font.name = "Georgia"
p.font.italic = True
p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
p.line_spacing = 1.5

# Attribution
attr = slide.shapes.add_textbox(Inches(1.5), Inches(5.0), Inches(10), Inches(0.6))
p = attr.text_frame.paragraphs[0]
p.text = "\u2014 Sarah Chen, CEO of TechForward Inc."
p.font.size = Pt(16)
p.font.name = "Calibri"
p.font.color.rgb = RGBColor(0x99, 0x00, 0x11)

# Bottom decorative line
bottom_line = slide.shapes.add_shape(1, Inches(1.5), Inches(5.8), Inches(4), Inches(0.03))
bottom_line.fill.solid()
bottom_line.fill.fore_color.rgb = RGBColor(0x44, 0x22, 0x28)
bottom_line.line.fill.background()
'''
            },
            {
                "name": "process_steps",
                "code": '''
# Numbered process steps with connecting elements - Sage Calm palette
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Sage green background
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = RGBColor(0x2D, 0x3A, 0x2D)

# Title
title = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11), Inches(0.9))
p = title.text_frame.paragraphs[0]
p.text = "Our Process"
p.font.size = Pt(40)
p.font.bold = True
p.font.name = "Georgia"
p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

steps = [
    {"num": "1", "title": "Discovery", "desc": "Deep analysis of your current state, pain points, and strategic objectives."},
    {"num": "2", "title": "Design", "desc": "Architecture and roadmap creation aligned with your business goals."},
    {"num": "3", "title": "Build", "desc": "Iterative development with continuous feedback and quality assurance."},
    {"num": "4", "title": "Launch", "desc": "Controlled deployment with monitoring, training, and ongoing optimization."},
]

step_width = 2.6
start_x = 0.8
spacing = 0.35

for i, step in enumerate(steps):
    x = start_x + i * (step_width + spacing)

    # Connecting line between steps (not after last)
    if i < len(steps) - 1:
        line_x = x + step_width
        conn = slide.shapes.add_shape(1, Inches(line_x), Inches(2.65), Inches(spacing), Inches(0.03))
        conn.fill.solid()
        conn.fill.fore_color.rgb = RGBColor(0x8F, 0xB9, 0x6A)
        conn.line.fill.background()

    # Number circle
    circle = slide.shapes.add_shape(9, Inches(x + step_width/2 - 0.35), Inches(2.0), Inches(0.7), Inches(0.7))
    circle.fill.solid()
    circle.fill.fore_color.rgb = RGBColor(0x8F, 0xB9, 0x6A)
    circle.line.fill.background()

    # Number text
    num = slide.shapes.add_textbox(Inches(x + step_width/2 - 0.35), Inches(2.08), Inches(0.7), Inches(0.6))
    p = num.text_frame.paragraphs[0]
    p.text = step["num"]
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.name = "Arial Black"
    p.font.color.rgb = RGBColor(0x2D, 0x3A, 0x2D)
    p.alignment = PP_ALIGN.CENTER

    # Step title
    st = slide.shapes.add_textbox(Inches(x), Inches(3.1), Inches(step_width), Inches(0.5))
    p = st.text_frame.paragraphs[0]
    p.text = step["title"]
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.name = "Calibri"
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.CENTER

    # Step description
    sd = slide.shapes.add_textbox(Inches(x + 0.15), Inches(3.7), Inches(step_width - 0.3), Inches(2.5))
    tf = sd.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = step["desc"]
    p.font.size = Pt(13)
    p.font.name = "Calibri"
    p.font.color.rgb = RGBColor(0xCC, 0xDD, 0xCC)
    p.line_spacing = 1.5
'''
            }
        ]
    },

    "design_reference": {
        "description": "Color palettes and font pairing reference for professional slide design",
        "when_to_use": "Reference when choosing colors and fonts for any slide",
        "is_text_reference": True,
        "examples": [
            {
                "name": "color_palettes",
                "text": """## Named Color Palettes

| Palette | Primary (BG) | Accent | Use Case |
|---------|-------------|--------|----------|
| Midnight Executive | #1E2761 | #408EC6 | Corporate, strategy |
| Teal Trust | #0A1A2A | #028090 | Trust, innovation |
| Forest & Moss | #2C5F2D | #97BC62 | Growth, sustainability |
| Berry & Cream | #ECE2D0 (light) | #6D2E46 | Warm, approachable data |
| Coral Energy | #1A1A2E | #FF6F61 / #00B4AA | Before/after, comparison |
| Ocean Gradient | #065A82 | #1B9AAA | Services, features |
| Warm Terracotta | #FAF5F0 (light) | #C46A4E | Product, showcase |
| Charcoal Minimal | #1C1C1E | #E8E8E8 | Hero, impact statements |
| Cherry Bold | #150E11 | #990011 | Quotes, emphasis |
| Sage Calm | #2D3A2D | #8FB96A | Process, steps |

## Usage Rules
- CRITICAL: Choose ONE palette for the ENTIRE presentation. All slides must share the same primary/accent colors.
- Dark backgrounds for emphasis and impact slides (60-70% of deck)
- Light backgrounds for data-heavy content slides (use a lighter tint of the same palette, not a different palette)
- Never use plain white (#FFFFFF) as a slide background
- Accent color should contrast strongly with background"""
            },
            {
                "name": "font_pairings",
                "text": """## Font Pairings

| Header Font | Body Font | Style |
|------------|-----------|-------|
| Georgia (serif) | Calibri (sans) | Classic professional |
| Arial Black | Arial | Bold modern |
| Calibri Bold | Calibri Light | Clean corporate |

## Typography Rules
- Titles: 36-44pt, bold, centered or left-aligned
- Subtitles: 18-22pt, regular or light weight
- Body text: 14-16pt, left-aligned, 1.4-1.6 line spacing
- Stats/numbers: 48-120pt, bold, accent color
- Labels: 12-16pt, uppercase with letter spacing for small labels
- Never center body paragraphs - left-align for readability

## Spacing Rules
- Minimum 0.5" margin from all slide edges
- 0.3-0.5" gap between elements
- Visual breathing room between title and content (0.5"+ gap)

## Anti-Patterns (AVOID)
- Plain bullet points on white background
- Default PowerPoint blue (#4472C4)
- Accent lines directly under titles
- Text-only slides without any visual elements
- More than 4 bullet points per slide
- Font sizes below 12pt"""
            }
        ]
    }
}


def get_examples(category: str = None) -> dict:
    """Get slide code examples.

    Args:
        category: One of "text_layout", "number_highlight", "grid_layout",
                  "image_text", "visual_emphasis", "design_reference",
                  or None for all

    Returns:
        Dict with examples for the category or all categories
    """
    if category and category in SLIDE_EXAMPLES:
        return {category: SLIDE_EXAMPLES[category]}
    return SLIDE_EXAMPLES


def get_all_categories() -> list:
    """Get list of available example categories."""
    return list(SLIDE_EXAMPLES.keys())
